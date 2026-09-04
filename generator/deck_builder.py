"""Moteur de génération des supports PPTX (thème + layouts réutilisables).

Chaque méthode publique de `Deck` produit une slide complète et renvoie l'objet
slide pour permettre des ajouts ponctuels. Le rendu est volontairement dépourvu
de placeholders PowerPoint : tout est dessiné pour garantir un rendu identique
sur Windows, macOS et LibreOffice.
"""

from __future__ import annotations

import re
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Palette
# --------------------------------------------------------------------------- #

INK = RGBColor(0x12, 0x18, 0x26)
INK_SOFT = RGBColor(0x4A, 0x54, 0x68)
INK_FAINT = RGBColor(0x8A, 0x93, 0xA6)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_ALT = RGBColor(0xF4, 0xF6, 0xFA)
NAVY = RGBColor(0x0B, 0x12, 0x20)
NAVY_SOFT = RGBColor(0x1B, 0x25, 0x3B)
RULE = RGBColor(0xDD, 0xE2, 0xEB)

CLAUDE = RGBColor(0xD9, 0x6B, 0x3C)   # orange Anthropic
COPILOT = RGBColor(0x2F, 0x81, 0xF7)  # bleu GitHub
GREEN = RGBColor(0x1F, 0x93, 0x5E)
AMBER = RGBColor(0xC2, 0x7A, 0x0A)
RED = RGBColor(0xC0, 0x39, 0x3C)
VIOLET = RGBColor(0x7A, 0x5A, 0xF8)

CODE_BG = RGBColor(0x14, 0x1B, 0x2D)
CODE_FG = RGBColor(0xE6, 0xEA, 0xF2)
CODE_COMMENT = RGBColor(0x7C, 0x8A, 0xA8)
CODE_ACCENT = RGBColor(0x7F, 0xD1, 0xAE)

FONT = "Segoe UI"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = SLIDE_W - 2 * MARGIN
TITLE_TOP = Inches(0.46)
BODY_TOP = Inches(1.62)
BODY_BOTTOM = Inches(6.86)
BODY_H = BODY_BOTTOM - BODY_TOP


@dataclass
class Badge:
    """Petite pastille de niveau / durée affichée en haut à droite."""

    label: str
    color: RGBColor


LEVEL_COLORS = {
    "INTERMÉDIAIRE": COPILOT,
    "EXPERT": VIOLET,
    "INTER. + EXPERT": GREEN,
    "ATELIER": CLAUDE,
    "DÉMO": AMBER,
}


# --------------------------------------------------------------------------- #
# Helpers bas niveau
# --------------------------------------------------------------------------- #


def _no_shadow(shape):
    shape.shadow.inherit = False


def _rect(slide, x, y, w, h, fill=None, line=None, radius=None, shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    _no_shadow(box)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(1)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        box.adjustments[0] = radius
    box.text_frame.word_wrap = True
    box.text_frame.text = ""
    return box


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def _style(run, size, color=INK, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


_MD_BOLD = re.compile(r"\*\*(.+?)\*\*", re.S)
_MD_CODE = re.compile(r"`([^`]+)`")
_MD_STRIP = re.compile(r"[*`]")


def _est_lines(text, width_in, size):
    """Estimation du nombre de lignes après retour à la ligne automatique."""
    plain = _MD_STRIP.sub("", text)
    char_w = 0.5 * size / 72.0
    per_line = max(12, int(width_in / char_w))
    return max(1, -(-len(plain) // per_line))


def _spread(texts, sizes, width_in, avail_emu, lo=6.0, hi=18.0):
    """Espacement inter-paragraphes qui répartit le bloc sur la hauteur utile.

    Évite les slides dont le contenu s'entasse en haut sur un tiers de la page.
    """
    lines_h = 0.0
    for text, size in zip(texts, sizes):
        lines_h += _est_lines(text, width_in, size) * size * 1.15 / 72.0
    leftover = avail_emu / 914400 - lines_h
    per_para = leftover / max(1, len(texts)) * 72.0
    return max(lo, min(hi, per_para * 0.72))


def _write_rich(paragraph, text, size, color=INK, accent=None, mono_color=None,
                bold=False):
    """Écrit `text` en interprétant **gras** et `code`, y compris imbriqués.

    Le jeton qui commence le plus tôt gagne, ce qui laisse les motifs glob
    (`**/*.tsx`) intacts lorsqu'ils sont écrits entre accents graves.
    """
    accent = accent or INK
    mono_color = mono_color or INK_SOFT
    pos = 0
    while pos < len(text):
        m_bold = _MD_BOLD.search(text, pos)
        m_code = _MD_CODE.search(text, pos)
        candidates = [m for m in (m_bold, m_code) if m]
        if not candidates:
            run = paragraph.add_run()
            run.text = text[pos:]
            _style(run, size, accent if bold else color, bold=bold)
            return
        match = min(candidates, key=lambda m: m.start())
        if match.start() > pos:
            run = paragraph.add_run()
            run.text = text[pos:match.start()]
            _style(run, size, accent if bold else color, bold=bold)
        if match is m_code:
            run = paragraph.add_run()
            run.text = match.group(1)
            _style(run, size - 0.5, accent if bold else mono_color,
                   bold=bold, font=FONT_MONO)
        else:
            _write_rich(paragraph, match.group(1), size, color, accent, mono_color,
                        bold=True)
        pos = match.end()


# --------------------------------------------------------------------------- #
# Deck
# --------------------------------------------------------------------------- #


class Deck:
    def __init__(self, footer: str, accent: RGBColor = CLAUDE):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.accent = accent
        self.footer = footer
        self._blank = self.prs.slide_layouts[6]
        self.section_index = 0

    # -- infrastructure ---------------------------------------------------- #

    def _new(self, dark=False, chrome=True):
        slide = self.prs.slides.add_slide(self._blank)
        bg = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY if dark else PAPER)
        bg.shadow.inherit = False
        if chrome:
            self._chrome(slide, dark)
        return slide

    def _chrome(self, slide, dark=False):
        y = SLIDE_H - Inches(0.42)
        _rect(slide, MARGIN, y, BODY_W, Emu(9525), fill=NAVY_SOFT if dark else RULE)
        _, tf = _textbox(slide, MARGIN, y + Inches(0.09), BODY_W * 0.75, Inches(0.26))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = self.footer
        _style(r, 9, INK_FAINT if not dark else RGBColor(0x6C, 0x78, 0x92))
        num = len(self.prs.slides._sldIdLst) if hasattr(self.prs.slides, "_sldIdLst") else 0
        _, tf2 = _textbox(slide, MARGIN + BODY_W * 0.8, y + Inches(0.09), BODY_W * 0.2, Inches(0.26))
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = str(num)
        _style(r2, 9, INK_FAINT if not dark else RGBColor(0x6C, 0x78, 0x92), bold=True)

    def _header(self, slide, title, kicker=None, badges=None, dark=False):
        """Bandeau titre + kicker + badges. Renvoie le bas du bandeau."""
        top = TITLE_TOP
        badges = badges or []
        title_w = BODY_W - (Inches(2.9) if badges else Inches(0))

        if kicker:
            _, tf = _textbox(slide, MARGIN, top, title_w, Inches(0.28))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = kicker.upper()
            _style(r, 10.5, self.accent, bold=True)
            r.font.name = FONT
            top += Inches(0.34)

        _, tf = _textbox(slide, MARGIN, top, title_w, Inches(0.9))
        p = tf.paragraphs[0]
        p.line_spacing = 0.95
        _write_rich(p, title, 27, PAPER if dark else INK, accent=self.accent)

        if badges:
            bx = MARGIN + BODY_W
            for badge in reversed(badges):
                w = Inches(0.42) + Inches(0.082) * len(badge.label)
                bx -= w
                box = _rect(
                    slide, bx, TITLE_TOP + Inches(0.04), w, Inches(0.3),
                    fill=badge.color, radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                )
                tf = box.text_frame
                tf.margin_left = tf.margin_right = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = badge.label
                _style(r, 9, PAPER, bold=True)
                bx -= Inches(0.1)

        rule_y = TITLE_TOP + (Inches(1.28) if kicker else Inches(0.98))
        _rect(slide, MARGIN, rule_y, Inches(1.5), Inches(0.035), fill=self.accent)
        return rule_y + Inches(0.28)

    @staticmethod
    def _notes(slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = text.strip()

    # -- layouts ----------------------------------------------------------- #

    def title_slide(self, eyebrow, title, subtitle, meta_lines, notes=None):
        slide = self._new(dark=True, chrome=False)
        _rect(slide, 0, 0, Inches(0.28), SLIDE_H, fill=self.accent)

        _, tf = _textbox(slide, Inches(1.35), Inches(1.55), Inches(10.4), Inches(0.3))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = eyebrow.upper()
        _style(r, 12, self.accent, bold=True)

        _, tf = _textbox(slide, Inches(1.35), Inches(2.05), Inches(10.6), Inches(1.9))
        p = tf.paragraphs[0]
        p.line_spacing = 0.94
        r = p.add_run()
        r.text = title
        _style(r, 46, PAPER, bold=True)

        _, tf = _textbox(slide, Inches(1.35), Inches(3.72), Inches(9.6), Inches(0.7))
        p = tf.paragraphs[0]
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = subtitle
        _style(r, 16.5, RGBColor(0xB6, 0xC0, 0xD4))

        _rect(slide, Inches(1.35), Inches(4.72), Inches(1.6), Inches(0.035), fill=self.accent)

        _, tf = _textbox(slide, Inches(1.35), Inches(5.14), Inches(10.4), Inches(1.4))
        for i, line in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            _write_rich(p, line, 12.5, RGBColor(0x93, 0x9F, 0xB8), accent=PAPER,
                        mono_color=CODE_ACCENT)
        self._notes(slide, notes)
        return slide

    def section(self, number, title, subtitle, bullets=None, notes=None):
        slide = self._new(dark=True, chrome=False)
        _rect(slide, 0, 0, Inches(0.28), SLIDE_H, fill=self.accent)
        _, tf = _textbox(slide, Inches(1.35), Inches(2.05), Inches(3.0), Inches(1.5))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = number
        _style(r, 80, NAVY_SOFT, bold=True)

        _, tf = _textbox(slide, Inches(1.35), Inches(3.25), Inches(10.4), Inches(1.2))
        p = tf.paragraphs[0]
        p.line_spacing = 0.96
        r = p.add_run()
        r.text = title
        _style(r, 34, PAPER, bold=True)

        _, tf = _textbox(slide, Inches(1.35), Inches(4.55), Inches(9.8), Inches(0.6))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        _style(r, 15, self.accent)

        if bullets:
            _, tf = _textbox(slide, Inches(1.35), Inches(5.3), Inches(10.2), Inches(1.5))
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(7)
                r = p.add_run()
                r.text = "— "
                _style(r, 12.5, NAVY_SOFT)
                _write_rich(p, b, 12.5, RGBColor(0x9E, 0xAA, 0xC2), accent=PAPER,
                            mono_color=CODE_ACCENT)
        self._notes(slide, notes)
        return slide

    def bullets(self, title, items, kicker=None, badges=None, notes=None, lead=None,
                size=15.0):
        """items: str, ou (str, niveau 0/1), niveau 1 = sous-point."""
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13.5, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.62)

        norm = [(it, 0) if isinstance(it, str) else it for it in items]
        gap_pt = _spread(
            [t for t, _ in norm],
            [size if lvl == 0 else size - 1.5 for _, lvl in norm],
            (BODY_W - Inches(0.3)) / 914400,
            BODY_BOTTOM - top,
            lo=7.0, hi=20.0,
        )
        _, tf = _textbox(slide, MARGIN, top, BODY_W, BODY_BOTTOM - top)
        for i, item in enumerate(items):
            text, lvl = (item, 0) if isinstance(item, str) else item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.12
            p.space_after = Pt(gap_pt if lvl == 0 else gap_pt * 0.45)
            if lvl:
                p.left_indent = Inches(0.52)
                p.first_line_indent = Inches(-0.22)
                r = p.add_run()
                r.text = "· "
                _style(r, size - 1.5, INK_FAINT)
                _write_rich(p, text, size - 1.5, INK_SOFT, accent=INK, mono_color=self.accent)
            else:
                p.left_indent = Inches(0.3)
                p.first_line_indent = Inches(-0.3)
                r = p.add_run()
                r.text = "▪  "
                _style(r, size - 2, self.accent, bold=True)
                _write_rich(p, text, size, INK, accent=self.accent, mono_color=INK_SOFT)
        self._notes(slide, notes)
        return slide

    def two_col(self, title, left, right, kicker=None, badges=None, notes=None,
                lead=None, size=13.5):
        """left/right : dict(title=..., color=..., items=[...])."""
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.6)

        gap = Inches(0.36)
        col_w = (BODY_W - gap) / 2
        for idx, col in enumerate((left, right)):
            x = MARGIN + idx * (col_w + gap)
            color = col.get("color", self.accent)
            h = BODY_BOTTOM - top
            _rect(slide, x, top, col_w, h, fill=PAPER_ALT, radius=0.04,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, x, top, col_w, Inches(0.045), fill=color)

            _, tf = _textbox(slide, x + Inches(0.3), top + Inches(0.28),
                             col_w - Inches(0.6), Inches(0.4))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = col["title"]
            _style(r, 15, color, bold=True)

            _, tf = _textbox(slide, x + Inches(0.3), top + Inches(0.82),
                             col_w - Inches(0.6), h - Inches(1.1))
            for i, item in enumerate(col["items"]):
                text, lvl = (item, 0) if isinstance(item, str) else item
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.line_spacing = 1.12
                p.space_after = Pt(8 if lvl == 0 else 4)
                p.left_indent = Inches(0.26 if lvl == 0 else 0.5)
                p.first_line_indent = Inches(-0.26)
                r = p.add_run()
                r.text = ("▪  " if lvl == 0 else "· ")
                _style(r, size - 2, color if lvl == 0 else INK_FAINT, bold=(lvl == 0))
                _write_rich(p, text, size - (0 if lvl == 0 else 1), INK if lvl == 0 else INK_SOFT,
                            accent=INK, mono_color=color)
        self._notes(slide, notes)
        return slide

    def code(self, title, code, kicker=None, badges=None, caption=None, notes=None,
             lead=None, size=None):
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.58)

        lines = code.strip("\n").split("\n")
        if size is None:
            n = len(lines)
            longest = max((len(l) for l in lines), default=0)
            size = 13.0 if n <= 12 else 11.5 if n <= 17 else 10.0 if n <= 22 else 8.8
            if longest > 92:
                size = min(size, 9.6)
            elif longest > 78:
                size = min(size, 11.0)

        cap_h = Inches(0.46) if caption else Inches(0)
        box_h = BODY_BOTTOM - top - cap_h
        _rect(slide, MARGIN, top, BODY_W, box_h, fill=CODE_BG, radius=0.03,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, MARGIN, top, Inches(0.06), box_h, fill=self.accent)

        _, tf = _textbox(slide, MARGIN + Inches(0.34), top + Inches(0.24),
                         BODY_W - Inches(0.6), box_h - Inches(0.4))
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.16
            p.space_after = Pt(0)
            stripped = line.strip()
            run = p.add_run()
            run.text = line if line else " "
            if stripped.startswith("#") or stripped.startswith("//"):
                _style(run, size, CODE_COMMENT, font=FONT_MONO, italic=True)
            elif stripped.startswith("$") or stripped.startswith(">"):
                _style(run, size, CODE_ACCENT, font=FONT_MONO, bold=True)
            else:
                _style(run, size, CODE_FG, font=FONT_MONO)

        if caption:
            _, tf = _textbox(slide, MARGIN, BODY_BOTTOM - Inches(0.36), BODY_W, Inches(0.36))
            p = tf.paragraphs[0]
            _write_rich(p, caption, 11.5, INK_SOFT, accent=self.accent, mono_color=self.accent)
        self._notes(slide, notes)
        return slide

    def table(self, title, headers, rows, kicker=None, badges=None, notes=None,
              lead=None, widths=None, size=11.5):
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.6)

        n_rows = len(rows) + 1
        avail = BODY_BOTTOM - top
        gf = slide.shapes.add_table(n_rows, len(headers), MARGIN, top, BODY_W, avail)
        tbl = gf.table
        tbl.first_row = False
        tbl.horz_banding = False

        if widths:
            total = sum(widths)
            for i, w in enumerate(widths):
                tbl.columns[i].width = Emu(int(BODY_W * w / total))

        head_h = Inches(0.44)
        tbl.rows[0].height = head_h
        body_h = max(Inches(0.34), int((avail - head_h) / max(1, len(rows))))
        for ri in range(1, n_rows):
            tbl.rows[ri].height = body_h

        for c, text in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.margin_left = cell.margin_right = Inches(0.14)
            cell.margin_top = cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = text
            _style(r, size, PAPER, bold=True)

        for ri, row in enumerate(rows, start=1):
            for ci, text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER if ri % 2 else PAPER_ALT
                cell.margin_left = cell.margin_right = Inches(0.14)
                cell.margin_top = cell.margin_bottom = Inches(0.05)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                p.line_spacing = 1.05
                _write_rich(p, str(text), size, INK_SOFT, accent=INK, mono_color=self.accent)
        self._notes(slide, notes)
        return slide

    def key_idea(self, statement, attribution=None, kicker=None, notes=None):
        slide = self._new(dark=True, chrome=False)
        _rect(slide, 0, 0, Inches(0.28), SLIDE_H, fill=self.accent)
        if kicker:
            _, tf = _textbox(slide, Inches(1.5), Inches(1.9), Inches(10.2), Inches(0.3))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = kicker.upper()
            _style(r, 11, self.accent, bold=True)

        _, tf = _textbox(slide, Inches(1.5), Inches(2.5), Inches(10.2), Inches(2.6))
        p = tf.paragraphs[0]
        p.line_spacing = 1.08
        _write_rich(p, statement, 28, RGBColor(0xE8, 0xEC, 0xF4), accent=self.accent,
                    mono_color=CODE_ACCENT)
        if attribution:
            _, tf = _textbox(slide, Inches(1.5), Inches(5.4), Inches(10.2), Inches(0.8))
            p = tf.paragraphs[0]
            p.line_spacing = 1.2
            _write_rich(p, attribution, 13, RGBColor(0x8E, 0x9B, 0xB5), accent=PAPER,
                        mono_color=CODE_ACCENT)
        self._notes(slide, notes)
        return slide

    def cards(self, title, cards, kicker=None, badges=None, notes=None, lead=None,
              cols=3):
        """cards: liste de dict(title, body, color=?, tag=?)."""
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.6)

        gap = Inches(0.28)
        rows = (len(cards) + cols - 1) // cols
        card_w = (BODY_W - gap * (cols - 1)) / cols
        card_h = (BODY_BOTTOM - top - gap * (rows - 1)) / rows
        # Une rangée unique de cartes courtes : éviter des blocs inutilement hauts
        max_h = Inches(3.5) if rows == 1 else Inches(2.6)
        if card_h > max_h:
            top += (card_h - max_h) / 2
            card_h = max_h
        for i, card in enumerate(cards):
            r, c = divmod(i, cols)
            x = MARGIN + c * (card_w + gap)
            y = top + r * (card_h + gap)
            color = card.get("color", self.accent)
            _rect(slide, x, y, card_w, card_h, fill=PAPER_ALT, radius=0.05,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, x, y, Inches(0.055), card_h, fill=color)

            _, tf = _textbox(slide, x + Inches(0.28), y + Inches(0.24),
                             card_w - Inches(0.5), Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.0
            r_ = p.add_run()
            r_.text = card["title"]
            _style(r_, 13.5, color, bold=True)

            _, tf = _textbox(slide, x + Inches(0.28), y + Inches(0.76),
                             card_w - Inches(0.5), card_h - Inches(1.0))
            body = card["body"] if isinstance(card["body"], list) else [card["body"]]
            for j, line in enumerate(body):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.line_spacing = 1.14
                p.space_after = Pt(6)
                _write_rich(p, line, 11.5, INK_SOFT, accent=INK, mono_color=color)
        self._notes(slide, notes)
        return slide

    def workshop(self, title, meta, steps, expected=None, trap=None, expert=None,
                 objective=None, notes=None, kicker="Atelier"):
        """Slide d'atelier : étapes à gauche, encadrés à droite."""
        slide = self._new()
        badges = [Badge(m, LEVEL_COLORS.get(m, self.accent)) for m in meta]
        top = self._header(slide, title, kicker, badges)

        if objective:
            _rect(slide, MARGIN, top, BODY_W, Inches(0.52), fill=PAPER_ALT, radius=0.08,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, MARGIN, top, Inches(0.05), Inches(0.52), fill=self.accent)
            _, tf = _textbox(slide, MARGIN + Inches(0.26), top + Inches(0.12),
                             BODY_W - Inches(0.5), Inches(0.36))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "Objectif  "
            _style(r, 11, self.accent, bold=True)
            _write_rich(p, objective, 12.5, INK, accent=INK, mono_color=INK_SOFT)
            top += Inches(0.72)

        side_w = Inches(4.05)
        left_w = BODY_W - side_w - Inches(0.34)
        norm = [(s, 0) if isinstance(s, str) else s for s in steps]
        gap_pt = _spread(
            [t for t, _ in norm],
            [13.0 if lvl == 0 else 11.5 for _, lvl in norm],
            (left_w - Inches(0.34)) / 914400,
            BODY_BOTTOM - top,
            lo=8.0, hi=22.0,
        )
        _, tf = _textbox(slide, MARGIN, top, left_w, BODY_BOTTOM - top)
        for i, step in enumerate(steps):
            text, lvl = (step, 0) if isinstance(step, str) else step
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.12
            p.space_after = Pt(gap_pt if lvl == 0 else gap_pt * 0.4)
            if lvl == 0:
                p.left_indent = Inches(0.34)
                p.first_line_indent = Inches(-0.34)
                num = sum(1 for s in steps[: i + 1] if isinstance(s, str) or s[1] == 0)
                r = p.add_run()
                r.text = f"{num}.  "
                _style(r, 12.5, self.accent, bold=True)
                _write_rich(p, text, 13, INK, accent=self.accent, mono_color=INK_SOFT)
            else:
                p.left_indent = Inches(0.58)
                p.first_line_indent = Inches(-0.2)
                r = p.add_run()
                r.text = "· "
                _style(r, 11.5, INK_FAINT)
                _write_rich(p, text, 11.5, INK_SOFT, accent=INK, mono_color=self.accent)

        panels = []
        if expected:
            panels.append(("Résultat attendu", expected, GREEN))
        if trap:
            panels.append(("Piège classique", trap, AMBER))
        if expert:
            panels.append(("Piste experte", expert, VIOLET))

        x = MARGIN + left_w + Inches(0.34)
        y = top
        avail = BODY_BOTTOM - top
        gap = Inches(0.16)
        if panels:
            # Hauteur proportionnelle au contenu : sinon un encadré long déborde
            # sur le suivant. On réduit le corps de texte jusqu'à ce que ça tienne.
            text_w = (side_w - Inches(0.46)) / 914400
            body_size = 11.0
            for candidate in (11.0, 10.5, 10.0, 9.5, 9.0):
                needed = []
                for _, body, _c in panels:
                    lines = body if isinstance(body, list) else [body]
                    n = sum(_est_lines(l, text_w, candidate) for l in lines)
                    needed.append(
                        Inches(0.78) + Inches(candidate * 1.16 / 72) * n
                        + Inches(0.07) * len(lines)
                    )
                total = sum(needed) + gap * (len(panels) - 1)
                body_size = candidate
                if total <= avail:
                    break
            extra = max(0, avail - (sum(needed) + gap * (len(panels) - 1)))
            heights = [int(n + extra / len(panels)) for n in needed]

            for (label, body, color), h in zip(panels, heights):
                _rect(slide, x, y, side_w, h, fill=PAPER_ALT, radius=0.06,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                _rect(slide, x, y, side_w, Inches(0.04), fill=color)
                _, tf = _textbox(slide, x + Inches(0.24), y + Inches(0.2),
                                 side_w - Inches(0.46), Inches(0.28))
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = label.upper()
                _style(r, 9.5, color, bold=True)
                _, tf = _textbox(slide, x + Inches(0.24), y + Inches(0.54),
                                 side_w - Inches(0.46), h - Inches(0.74))
                body = body if isinstance(body, list) else [body]
                for j, line in enumerate(body):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.line_spacing = 1.13
                    p.space_after = Pt(5)
                    _write_rich(p, line, body_size, INK_SOFT, accent=INK, mono_color=color)
                y += h + gap
        self._notes(slide, notes)
        return slide

    def agenda(self, title, rows, kicker=None, notes=None, highlight=None):
        """rows: (horaire, titre, type) — type: 'theorie' | 'atelier' | 'pause'."""
        slide = self._new()
        top = self._header(slide, title, kicker)
        colors = {"theorie": COPILOT, "atelier": CLAUDE, "pause": INK_FAINT,
                  "capstone": VIOLET}
        h = (BODY_BOTTOM - top) / len(rows)
        for i, (hour, label, kind) in enumerate(rows):
            y = top + i * h
            color = colors.get(kind, self.accent)
            is_hl = highlight is not None and i in highlight
            _rect(slide, MARGIN, y, BODY_W, h - Inches(0.045),
                  fill=PAPER_ALT if (i % 2 == 0 or is_hl) else PAPER, radius=0.06,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _rect(slide, MARGIN, y, Inches(0.05), h - Inches(0.045), fill=color)
            _, tf = _textbox(slide, MARGIN + Inches(0.24), y + Inches(0.07),
                             Inches(1.5), h - Inches(0.12))
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = hour
            _style(r, 11.5, color, bold=True, font=FONT_MONO)
            _, tf = _textbox(slide, MARGIN + Inches(1.85), y + Inches(0.07),
                             BODY_W - Inches(2.2), h - Inches(0.12))
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.line_spacing = 1.05
            _write_rich(p, label, 12.5, INK, accent=color, mono_color=INK_SOFT)
        self._notes(slide, notes)
        return slide

    def checklist(self, title, items, kicker=None, badges=None, notes=None, lead=None,
                  cols=2):
        slide = self._new()
        top = self._header(slide, title, kicker, badges)
        if lead:
            _, tf = _textbox(slide, MARGIN, top, BODY_W, Inches(0.5))
            p = tf.paragraphs[0]
            p.line_spacing = 1.15
            _write_rich(p, lead, 13, INK_SOFT, accent=INK, mono_color=self.accent)
            top += Inches(0.6)

        per_col = (len(items) + cols - 1) // cols
        gap = Inches(0.4)
        col_w = (BODY_W - gap * (cols - 1)) / cols
        for c in range(cols):
            chunk = items[c * per_col : (c + 1) * per_col]
            if not chunk:
                continue
            x = MARGIN + c * (col_w + gap)
            gap_pt = _spread(chunk, [12.5] * len(chunk),
                             (col_w - Inches(0.34)) / 914400,
                             BODY_BOTTOM - top, lo=8.0, hi=22.0)
            _, tf = _textbox(slide, x, top, col_w, BODY_BOTTOM - top)
            for i, item in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.line_spacing = 1.14
                p.space_after = Pt(gap_pt)
                p.left_indent = Inches(0.34)
                p.first_line_indent = Inches(-0.34)
                r = p.add_run()
                r.text = "☐  "
                _style(r, 13, self.accent, bold=True)
                _write_rich(p, item, 12.5, INK_SOFT, accent=INK, mono_color=self.accent)
        self._notes(slide, notes)
        return slide

    def save(self, path):
        self.prs.save(path)
        return path
